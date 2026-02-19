"""
Celery 任务
"""
from celery import Task
from datetime import datetime
import os
import tempfile
import structlog
import oss2

from app.celery_app import celery_app
from app.database import SessionLocal
from app.models import Task as TaskModel
from app.config import settings

logger = structlog.get_logger()


class DatabaseTask(Task):
    """带数据库会话的任务基类"""

    _db = None

    @property
    def db(self):
        if self._db is None:
            self._db = SessionLocal()
        return self._db

    def after_return(self, *args, **kwargs):
        if self._db is not None:
            self._db.close()


@celery_app.task(bind=True, base=DatabaseTask, name="app.tasks.convert_pdf_task")
def convert_pdf_task(self, task_id: str, oss_key: str, task_type: str):
    """
    PDF 转换任务

    Args:
        task_id: 任务 ID
        oss_key: OSS 文件路径
        task_type: 转换类型 (pdf2word, pdf2excel, pdf2ppt)
    """
    logger.info("Starting PDF conversion", task_id=task_id, task_type=task_type)

    # 更新任务状态为处理中
    task = self.db.query(TaskModel).filter(TaskModel.task_id == task_id).first()
    if not task:
        logger.error("Task not found", task_id=task_id)
        return

    task.status = "processing"
    self.db.commit()

    try:
        # 初始化 OSS 客户端
        auth = oss2.Auth(settings.OSS_ACCESS_KEY, settings.OSS_SECRET_KEY)
        bucket = oss2.Bucket(auth, f"https://{settings.OSS_ENDPOINT}", settings.OSS_BUCKET)

        # 下载 PDF 文件到临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, "input.pdf")
            output_path = os.path.join(temp_dir, f"output.{get_output_extension(task_type)}")

            logger.info("Downloading PDF from OSS", oss_key=oss_key)
            bucket.get_object_to_file(oss_key, input_path)

            # 执行转换
            logger.info("Converting PDF", task_type=task_type)
            convert_pdf(input_path, output_path, task_type)

            # 上传结果到 OSS
            result_key = f"results/{task_id}.{get_output_extension(task_type)}"
            logger.info("Uploading result to OSS", result_key=result_key)
            bucket.put_object_from_file(result_key, output_path)

            # 更新任务状态
            task.status = "completed"
            task.oss_key_result = result_key
            task.completed_at = datetime.now()
            self.db.commit()

            logger.info("PDF conversion completed", task_id=task_id)

    except Exception as e:
        logger.error("PDF conversion failed", task_id=task_id, error=str(e))

        task.status = "failed"
        task.error_msg = str(e)
        task.completed_at = datetime.now()
        self.db.commit()

        raise


def convert_pdf(input_path: str, output_path: str, task_type: str):
    """
    执行 PDF 转换

    这是一个简化的实现，实际应该调用具体的转换库
    """
    if task_type == "pdf2word":
        # 使用 pdf2docx 转换
        try:
            from pdf2docx import Converter

            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
        except ImportError:
            # 降级方案：使用 PyPDF2 提取文本
            import PyPDF2
            from docx import Document

            with open(input_path, "rb") as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                doc = Document()

                for page in pdf_reader.pages:
                    text = page.extract_text()
                    doc.add_paragraph(text)

                doc.save(output_path)

    elif task_type == "pdf2excel":
        # 使用 pdfplumber 提取表格
        import pdfplumber
        import pandas as pd

        with pdfplumber.open(input_path) as pdf:
            tables = []
            for page in pdf.pages:
                page_tables = page.extract_tables()
                tables.extend(page_tables)

            # 将所有表格写入 Excel
            if tables:
                writer = pd.ExcelWriter(output_path, engine="openpyxl")
                for i, table in enumerate(tables):
                    df = pd.DataFrame(table[1:], columns=table[0])
                    df.to_excel(writer, sheet_name=f"Sheet{i+1}", index=False)
                writer.close()

    elif task_type == "pdf2ppt":
        # 简化实现：将 PDF 转为图片，然后插入 PPT
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches

        images = convert_from_path(input_path)
        prs = Presentation()

        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            img_path = os.path.join(os.path.dirname(output_path), f"page_{len(prs.slides)}.png")
            img.save(img_path, "PNG")

            # 添加图片到幻灯片
            left = top = Inches(0)
            slide.shapes.add_picture(img_path, left, top, width=prs.slide_width)

            os.remove(img_path)

        prs.save(output_path)

    else:
        raise ValueError(f"Unsupported task type: {task_type}")


def get_output_extension(task_type: str) -> str:
    """获取输出文件扩展名"""
    extensions = {"pdf2word": "docx", "pdf2excel": "xlsx", "pdf2ppt": "pptx", "merge": "pdf", "split": "pdf"}

    return extensions.get(task_type, "bin")


@celery_app.task(name="app.tasks.cleanup_expired_files")
def cleanup_expired_files():
    """清理过期文件定时任务"""
    logger.info("Starting cleanup expired files")

    db = SessionLocal()

    try:
        # 查询过期任务
        expired_tasks = (
            db.query(TaskModel)
            .filter(TaskModel.status == "completed", TaskModel.expire_at < datetime.now())
            .all()
        )

        logger.info("Found expired tasks", count=len(expired_tasks))

        # 初始化 OSS 客户端
        auth = oss2.Auth(settings.OSS_ACCESS_KEY, settings.OSS_SECRET_KEY)
        bucket = oss2.Bucket(auth, f"https://{settings.OSS_ENDPOINT}", settings.OSS_BUCKET)

        deleted_count = 0

        for task in expired_tasks:
            try:
                # 删除源文件
                if task.oss_key_source:
                    bucket.delete_object(task.oss_key_source)

                # 删除结果文件
                if task.oss_key_result:
                    bucket.delete_object(task.oss_key_result)

                # 更新任务状态
                task.status = "expired"
                deleted_count += 1

            except Exception as e:
                logger.error("Failed to delete file", task_id=task.task_id, error=str(e))

        db.commit()
        logger.info("Cleanup completed", deleted_count=deleted_count)

    finally:
        db.close()


@celery_app.task(name="app.tasks.generate_hourly_stats")
def generate_hourly_stats():
    """生成每小时统计数据"""
    logger.info("Generating hourly stats")

    db = SessionLocal()

    try:
        from sqlalchemy import func

        # 统计最近 1 小时的任务
        one_hour_ago = datetime.now() - timedelta(hours=1)

        stats = (
            db.query(
                func.count(TaskModel.task_id).label("total"),
                func.sum(func.case((TaskModel.status == "completed", 1), else_=0)).label("completed"),
                func.sum(func.case((TaskModel.status == "failed", 1), else_=0)).label("failed"),
            )
            .filter(TaskModel.created_at >= one_hour_ago)
            .first()
        )

        logger.info(
            "Hourly stats generated",
            total=stats.total,
            completed=stats.completed,
            failed=stats.failed,
        )

    finally:
        db.close()
